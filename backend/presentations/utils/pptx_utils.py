# your_app/utils/pptx_utils.py
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from io import BytesIO
import requests
from typing import Optional, List


def _download_image_bytes(url: str) -> Optional[BytesIO]:
    print("⬇️ [_download_image_bytes] Скачиваю изображение:", url)

    try:
        resp = requests.get(url, timeout=15)
        print("⬇️ [_download_image_bytes] Статус ответа:", resp.status_code)
        resp.raise_for_status()
        print("⬇️ [_download_image_bytes] Размер полученного контента:", len(resp.content))
        b = BytesIO(resp.content)
        b.seek(0)
        return b
    except Exception as e:
        print("❌ [_download_image_bytes] Ошибка скачивания:", e)
        return None



def _remove_shape(slide, shape):
    # аккуратно удаляем shape из XML
    el = shape._element
    el.getparent().remove(el)


def _replace_pictures_on_slide_with_image(slide, image_stream):
    print(f"🖼️ [_replace_pictures_on_slide_with_image] Обрабатываю слайд")
    shapes = list(slide.shapes)
    found = False

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            found = True
            print("🔁 Найден picture shape → заменяю")
            print("   Позиция:", shape.left, shape.top, shape.width, shape.height)

            left, top, width, height = shape.left, shape.top, shape.width, shape.height

            try:
                slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
                print("   ✅ add_picture ОК")
            except Exception as e:
                print("   ❌ add_picture Ошибка:", e)

            try:
                _remove_shape(slide, shape)
                print("   🗑 Старое изображение удалено")
            except Exception as e:
                print("   ❌ Ошибка удаления старой картинки:", e)

            image_stream.seek(0)

    if not found:
        print("⚠️ Нет PICTURE shapes на этом слайде")


def _remove_picture_shapes_on_slide(slide):
    """Удаляет все picture shapes на слайде (placeholder'ы)."""
    shapes = list(slide.shapes)
    removed = 0
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                _remove_shape(slide, shape)
                removed += 1
        except Exception as e:
            print("❌ Ошибка при удалении picture shape:", e)
    if removed:
        print(f"🗑 Удалено picture shapes: {removed}")
    else:
        print("⚠️ Нет picture shapes для удаления")



def _set_text_preserve_run_style(shape, new_text: str):
    """
    Сохраняет стиль (цвет/шрифт/size/bold) первого run'а и заменяет текст в первом run'е.
    Если есть дополнительные run'ы — их текст обнуляется.
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    if not tf.paragraphs:
        # создаём параграф и run
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = new_text
        return

    p = tf.paragraphs[0]
    if p.runs:
        first_run = p.runs[0]
        # Записываем новый текст в первый run (сохраняем стиль)
        first_run.text = new_text
        # очищаем последующие run'ы, чтобы не было наслоения
        for run in p.runs[1:]:
            run.text = ""
    else:
        # нет runs — просто задаём paragraph.text (создаст run с дефолт. стилем)
        p.text = new_text


def fill_pptx_template(template_file_path: str, slides_data: dict, image_urls: Optional[List[str]] = None) -> BytesIO:
    """
    Заполняет PPTX-шаблон данными, заменяет изображения на 2 и 4 слайдах (если image_url задан)
    и возвращает BytesIO с готовым PPTX.
    - template_file_path: путь к .pptx шаблону
    - slides_data: {"slides": [{...}, ...]}
    - image_url: если задан — изображение будет вставлено в 2-й и 4-й слайды (если есть picture shapes)
    """
    prs = Presentation(template_file_path)
    slides_list = slides_data.get("slides", [])

    # Подготовим image streams по слайдам (2-й и 4-й -> индексы 1 и 3)
    # image_urls может быть None, пустым списком, либо содержать 1 или 2 URL'а
    slide_image_streams = {1: None, 3: None}
    if image_urls:
        if len(image_urls) > 0 and image_urls[0]:
            slide_image_streams[1] = _download_image_bytes(image_urls[0])
        if len(image_urls) > 1 and image_urls[1]:
            slide_image_streams[3] = _download_image_bytes(image_urls[1])

    for idx, slide in enumerate(prs.slides):
        if idx < len(slides_list):
            slide_data = slides_list[idx]

            # Заголовок
            if slide.shapes.title:
                _set_text_preserve_run_style(slide.shapes.title, slide_data.get("title", ""))

            # Основной текст — ищем текстовые shapes кроме заголовка
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape == slide.shapes.title:
                    continue
                _set_text_preserve_run_style(shape, slide_data.get("text", ""))

        # Обработка картинок для 2-го и 4-го слайдов (индексы 1 и 3)
        if idx in (1, 3):
            stream = slide_image_streams.get(idx)
            if stream:
                try:
                    stream.seek(0)
                    _replace_pictures_on_slide_with_image(slide, stream)
                except Exception as e:
                    print("❌ Ошибка вставки картинки на слайд", idx + 1, e)
                    _remove_picture_shapes_on_slide(slide)
            else:
                # если нет изображения для этого слайда — удаляем плейсхолдеры
                _remove_picture_shapes_on_slide(slide)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

